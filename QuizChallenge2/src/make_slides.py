"""Build the 10-slide pptx programmatically.

Run: python -m src.make_slides
Outputs: slides/QuizChallenge2.pptx
"""
from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).parent.parent
PLOTS = ROOT / "results" / "plots"
OUT = ROOT / "slides" / "QuizChallenge2.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)

WIDE = Inches(13.333)
TALL = Inches(7.5)

prs = Presentation()
prs.slide_width = WIDE
prs.slide_height = TALL

BLANK = prs.slide_layouts[6]
NAVY = RGBColor(0x1A, 0x33, 0x5C)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)
GREY = RGBColor(0x55, 0x55, 0x55)


def add_textbox(slide, x, y, w, h, text, *, size=18, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    if isinstance(text, str):
        para = tf.paragraphs[0]
        para.alignment = align
        run = para.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    else:  # list of (text, opts)
        for i, item in enumerate(text):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.alignment = item.get("align", PP_ALIGN.LEFT)
            run = para.add_run()
            run.text = item["text"]
            run.font.size = Pt(item.get("size", size))
            run.font.bold = item.get("bold", bold)
            run.font.color.rgb = item.get("color", color)
    return tb


def add_image(slide, path, x, y, w=None, h=None):
    if path.exists():
        slide.shapes.add_picture(str(path), x, y, width=w, height=h)


def title_slide(prs):
    s = prs.slides.add_slide(BLANK)
    add_textbox(
        s,
        Inches(0.7), Inches(2.4), Inches(12), Inches(1.2),
        "MelodyLeak",
        size=64, bold=True, color=NAVY,
    )
    add_textbox(
        s,
        Inches(0.7), Inches(3.5), Inches(12), Inches(0.8),
        "Foundation-Model Music as a Covert Speech Channel",
        size=28, color=ACCENT,
    )
    add_textbox(
        s,
        Inches(0.7), Inches(5.2), Inches(12), Inches(1.2),
        [
            {"text": "Murali Krishna Ediga", "size": 22, "bold": True, "color": NAVY},
            {"text": "CS 5542 — Big Data Analytics & AI · Spring 2026", "size": 18, "color": GREY},
            {"text": "University of Missouri — Kansas City · Quiz Challenge 2", "size": 16, "color": GREY},
        ],
    )


def header(slide, title):
    add_textbox(
        slide,
        Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.7),
        title,
        size=30, bold=True, color=NAVY,
    )


def slide_2(prs):
    s = prs.slides.add_slide(BLANK)
    header(s, "Problem & motivation")
    add_textbox(
        s, Inches(0.7), Inches(1.2), Inches(12), Inches(0.7),
        "Generative audio foundation models are now production-grade in BOTH directions:",
        size=22, color=NAVY,
    )
    add_textbox(
        s, Inches(1.0), Inches(2.0), Inches(11), Inches(2.0),
        [
            {"text": "•  Generate music + speech: MusicGen, SpeechT5, AudioLDM2, Bark", "size": 20},
            {"text": "•  Transcribe speech: Whisper, wav2vec2", "size": 20},
        ],
    )
    add_textbox(
        s, Inches(0.7), Inches(4.2), Inches(12), Inches(0.8),
        "Question",
        size=22, bold=True, color=ACCENT,
    )
    add_textbox(
        s, Inches(0.7), Inches(4.9), Inches(12), Inches(1.2),
        "Can a single pipeline use both ends to smuggle a hidden message inside generated music?",
        size=22, color=NAVY,
    )
    add_textbox(
        s, Inches(0.7), Inches(6.0), Inches(12), Inches(1.2),
        "Why it matters: as generative-AI audio proliferates online, content provenance and covert-channel detection become open security problems. Demonstrating channels of this kind is a prerequisite to building defences.",
        size=16, color=GREY,
    )


def slide_3(prs):
    s = prs.slides.add_slide(BLANK)
    header(s, "Models used  —  HuggingFace, inference only, no fine-tuning")
    rows = [
        ("Role", "HuggingFace model", "Params", "License"),
        ("TTS  (Vector A encoder)", "microsoft/speecht5_tts + speecht5_hifigan", "144 M", "MIT"),
        ("Carrier  (both vectors)", "facebook/musicgen-small", "300 M", "CC-BY-NC-4.0"),
        ("Decoder  (Vector A)", "openai/whisper-base", "74 M", "MIT"),
        ("Decoder + perceptual  (Vector B)", "laion/clap-htsat-unfused", "153 M", "CC0"),
    ]
    table = s.shapes.add_table(len(rows), 4, Inches(0.5), Inches(1.4), Inches(12.3), Inches(3.6)).table
    widths = [Inches(3.4), Inches(4.6), Inches(1.4), Inches(2.9)]
    for i, w in enumerate(widths):
        table.columns[i].width = w
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            tf = cell.text_frame
            tf.text = ""
            para = tf.paragraphs[0]
            run = para.add_run()
            run.text = val
            run.font.size = Pt(16 if r else 18)
            run.font.bold = (r == 0)
            run.font.color.rgb = NAVY if r else GREY if c == 2 else NAVY
    add_textbox(
        s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.7),
        "Footprint: ~3.4 GB VRAM total on float32. Single RTX 5080 runs everything live.",
        size=16, color=ACCENT,
    )
    add_textbox(
        s, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.7),
        "All training was done by the model authors. No data of mine touches model weights.",
        size=14, color=GREY,
    )


def slide_4(prs):
    s = prs.slides.add_slide(BLANK)
    header(s, "Two-vector design")
    add_textbox(
        s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.0),
        "Same secret routed through two complementary channels that share a MusicGen carrier.",
        size=18, color=GREY,
    )
    # Vector A box
    add_textbox(
        s, Inches(0.5), Inches(2.5), Inches(6), Inches(0.6),
        "Vector A — semantic", size=22, bold=True, color=ACCENT,
    )
    add_textbox(
        s, Inches(0.5), Inches(3.1), Inches(6), Inches(3.2),
        [
            {"text": "secret text", "size": 16, "bold": True},
            {"text": "↓", "size": 18},
            {"text": "SpeechT5 → vocals (16 kHz)", "size": 16},
            {"text": "MusicGen → carrier (32 kHz, resampled)", "size": 16},
            {"text": "↓ mix at gain α", "size": 16},
            {"text": "↓", "size": 18},
            {"text": "Whisper-base → recovered text", "size": 16, "bold": True},
        ],
    )
    add_textbox(
        s, Inches(0.5), Inches(6.3), Inches(6), Inches(1.0),
        "high bandwidth (~5–10 bps)  ·  fragile to gain attenuation  ·  audible to humans",
        size=14, color=GREY,
    )
    # Vector B box
    add_textbox(
        s, Inches(7.0), Inches(2.5), Inches(6), Inches(0.6),
        "Vector B — categorical", size=22, bold=True, color=ACCENT,
    )
    add_textbox(
        s, Inches(7.0), Inches(3.1), Inches(6), Inches(3.2),
        [
            {"text": "secret bits (2 bits per clip)", "size": 16, "bold": True},
            {"text": "↓", "size": 18},
            {"text": "MusicGen with prompt from a 4-class codebook", "size": 16},
            {"text": "(mood × tempo)", "size": 14, "color": GREY},
            {"text": "↓", "size": 18},
            {"text": "CLAP cosine to anchor prompts → argmax", "size": 16},
            {"text": "→ recovered bits", "size": 16, "bold": True},
        ],
    )
    add_textbox(
        s, Inches(7.0), Inches(6.3), Inches(6), Inches(1.0),
        "low bandwidth (~0.33 bps)  ·  robust to MP3 / time-stretch  ·  carrier looks like normal music",
        size=14, color=GREY,
    )


def slide_5(prs):
    s = prs.slides.add_slide(BLANK)
    header(s, "Prompt / input engineering")
    add_textbox(
        s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.6),
        "Vector A — two mix variants compared", size=22, bold=True, color=ACCENT,
    )
    add_textbox(
        s, Inches(0.7), Inches(1.9), Inches(12.3), Inches(2.0),
        [
            {"text": "•  flat: full-band TTS vocals + music   (baseline)", "size": 18},
            {"text": "•  bandpass: TTS vocals filtered to 300–3400 Hz before mixing  (improved hypothesis)", "size": 18},
            {"text": "Hypothesis: telephone-band vocals match Whisper's training distribution → lower WER at same loudness.", "size": 16, "color": GREY},
        ],
    )
    add_textbox(
        s, Inches(0.5), Inches(4.1), Inches(12.3), Inches(0.6),
        "Vector B — 4-prompt codebook over mood × tempo", size=22, bold=True, color=ACCENT,
    )
    rows = [
        ("bits", "prompt"),
        ("00", "happy upbeat acoustic guitar pop, 140 bpm, major"),
        ("01", "happy warm soft acoustic guitar, 60 bpm, major"),
        ("10", "sad melancholic piano, 140 bpm, minor"),
        ("11", "sad melancholic piano, 60 bpm, minor"),
    ]
    table = s.shapes.add_table(len(rows), 2, Inches(0.7), Inches(4.8), Inches(12), Inches(2.4)).table
    table.columns[0].width = Inches(1.3)
    table.columns[1].width = Inches(10.7)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            tf = cell.text_frame
            tf.text = ""
            para = tf.paragraphs[0]
            run = para.add_run()
            run.text = val
            run.font.size = Pt(16)
            run.font.bold = (r == 0)
            run.font.color.rgb = NAVY


def slide_plot(prs, title, plot_name, caption_lines):
    s = prs.slides.add_slide(BLANK)
    header(s, title)
    add_image(s, PLOTS / plot_name, Inches(0.6), Inches(1.2), w=Inches(8.4))
    add_textbox(
        s, Inches(9.3), Inches(1.4), Inches(3.7), Inches(5.5),
        caption_lines,
    )
    return s


def slide_6(prs):
    return slide_plot(
        prs,
        "Vector A results — WER vs α",
        "01_wer_vs_alpha.png",
        [
            {"text": "Findings", "size": 20, "bold": True, "color": ACCENT},
            {"text": "•  Ambient is the most permissive carrier — WER ≈ 0.18 across the entire α range (-6 to -30 dB).", "size": 14, "color": NAVY},
            {"text": "•  Classical is the worst — WER explodes past α=-18 dB; string formants compete with vocals in-band.", "size": 14, "color": NAVY},
            {"text": "•  Lo-fi sits between — usable until α=-30 dB.", "size": 14, "color": NAVY},
            {"text": "Insight", "size": 18, "bold": True, "color": ACCENT},
            {"text": "Carrier choice matters as much as gain level.", "size": 14, "color": NAVY},
            {"text": "n=45 runs (3 prompts × 5 α × 3 secrets)", "size": 12, "color": GREY},
        ],
    )


def slide_7(prs):
    return slide_plot(
        prs,
        "Vector A — bandpass surprise (counter-hypothesis result)",
        "02_flat_vs_bandpass.png",
        [
            {"text": "Hypothesis", "size": 18, "bold": True, "color": ACCENT},
            {"text": "Bandpass-shaping vocals (300–3400 Hz, telephone band) should reduce WER at the same loudness.", "size": 13, "color": NAVY},
            {"text": "Result", "size": 18, "bold": True, "color": ACCENT},
            {"text": "Bandpass HURTS — classical (0.51 → 1.78), lo-fi (0.16 → 0.32). Only ambient is roughly neutral.", "size": 13, "color": NAVY},
            {"text": "Interpretation", "size": 18, "bold": True, "color": ACCENT},
            {"text": "Whisper relies on out-of-band fricatives (sibilants /s/, /sh/ above 3.4 kHz) that the filter strips.", "size": 13, "color": NAVY},
            {"text": "The 'improvement' backfires.", "size": 13, "color": NAVY, "bold": True},
        ],
    )


def slide_8(prs):
    return slide_plot(
        prs,
        "Vector B — BER under codec / noise / time-stretch",
        "03_ber_vs_distortion.png",
        [
            {"text": "Headline", "size": 18, "bold": True, "color": ACCENT},
            {"text": "BER flat at 0.20 through MP3 (320/192/128 kbps), time-stretch ±5%, AWGN 30dB. Only AWGN 10dB pushes BER higher.", "size": 13, "color": NAVY},
            {"text": "Why 0.20?", "size": 18, "bold": True, "color": ACCENT},
            {"text": "Confusion matrix (clean): CLAP perfectly discriminates mood (happy vs sad, 0% conf) but is tempo-blind — all slow prompts get classified as fast.", "size": 13, "color": NAVY},
            {"text": "Implication", "size": 18, "bold": True, "color": ACCENT},
            {"text": "A 1-bit mood-only codebook would give BER = 0 — capacity halves, robustness jumps.", "size": 13, "color": NAVY, "bold": True},
            {"text": "n=140 decodes (20 clips × 7 distortions)", "size": 11, "color": GREY},
        ],
    )


def slide_9(prs):
    return slide_plot(
        prs,
        "Cross-evaluation — dual failure modes",
        "04_cross_eval_scatter.png",
        [
            {"text": "Vector A (red)", "size": 18, "bold": True, "color": ACCENT},
            {"text": "~50 bps speech, WER from 0.17 (α=-6) to 0.83 (α=-30). Fragile to gain attenuation.", "size": 13, "color": NAVY},
            {"text": "Vector B (blue)", "size": 18, "bold": True, "color": ACCENT},
            {"text": "~0.33 bps, BER 0.20 across all codec/time-stretch conditions. Fragile to noise but not codec.", "size": 13, "color": NAVY},
            {"text": "Insight", "size": 18, "bold": True, "color": ACCENT},
            {"text": "The two vectors fail in opposite directions. A defence that catches one does NOT catch the other. Defending generative-AI audio covert channels requires more than one detector.", "size": 13, "color": NAVY, "bold": True},
        ],
    )


def slide_10(prs):
    s = prs.slides.add_slide(BLANK)
    header(s, "Limitations  ·  Ethics  ·  Future work  ·  AI disclosure")
    add_textbox(
        s, Inches(0.5), Inches(1.2), Inches(6.2), Inches(0.6),
        "Limitations", size=20, bold=True, color=ACCENT,
    )
    add_textbox(
        s, Inches(0.7), Inches(1.85), Inches(6), Inches(2.4),
        [
            {"text": "•  No human listening test (n=0) — flagged as next-step", "size": 13},
            {"text": "•  English-only secrets, single TTS speaker (CMU Arctic xvector idx 7306)", "size": 13},
            {"text": "•  Single carrier model (MusicGen-small)", "size": 13},
            {"text": "•  No active steganalysis tested", "size": 13},
            {"text": "•  3 secrets × 3 prompts × 5 α — small N", "size": 13},
        ],
    )
    add_textbox(
        s, Inches(0.5), Inches(4.3), Inches(6.2), Inches(0.6),
        "Ethics", size=20, bold=True, color=ACCENT,
    )
    add_textbox(
        s, Inches(0.7), Inches(4.95), Inches(6), Inches(2.4),
        [
            {"text": "•  Demonstration of a covert channel as a prerequisite for DEFENSIVE research.", "size": 13},
            {"text": "•  Vector A is detectable by attentive listeners — NOT perceptually-secure stego.", "size": 13},
            {"text": "•  Vector B leaks via prompt-distribution biases over a session.", "size": 13},
            {"text": "•  Framed as measurement of attack surface, not deployment guide.", "size": 13},
        ],
    )
    add_textbox(
        s, Inches(7.0), Inches(1.2), Inches(6), Inches(0.6),
        "Future work", size=20, bold=True, color=ACCENT,
    )
    add_textbox(
        s, Inches(7.2), Inches(1.85), Inches(5.8), Inches(2.4),
        [
            {"text": "1.  Adversarial perturbation of MusicGen output to fool Whisper while preserving human listening.", "size": 13},
            {"text": "2.  Mood-only 1-bit Vector B + redundancy coding for near-zero BER.", "size": 13},
            {"text": "3.  Detection: train a classifier on (clean MusicGen) vs (vocal-mixed MusicGen).", "size": 13},
        ],
    )
    add_textbox(
        s, Inches(7.0), Inches(4.3), Inches(6), Inches(0.6),
        "AI tools disclosure", size=20, bold=True, color=ACCENT,
    )
    add_textbox(
        s, Inches(7.2), Inches(4.95), Inches(5.8), Inches(2.4),
        [
            {"text": "•  Anthropic Claude (Opus 4.7) — adversarial design review, code scaffolding, slide outline. Specific decisions logged in AI_DISCLOSURE.md.", "size": 13},
            {"text": "•  HuggingFace pretrained models — inference only.", "size": 13},
            {"text": "•  No model trained or fine-tuned by me.", "size": 13},
        ],
    )
    add_textbox(
        s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
        "Repo: github.com/muralikrish9/CS5542 — QuizChallenge2/  ·  Code, eval CSVs, plots, audio samples",
        size=14, color=NAVY, bold=True,
    )


def main():
    title_slide(prs)
    slide_2(prs)
    slide_3(prs)
    slide_4(prs)
    slide_5(prs)
    slide_6(prs)
    slide_7(prs)
    slide_8(prs)
    slide_9(prs)
    slide_10(prs)
    prs.save(OUT)
    print(f"saved {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
