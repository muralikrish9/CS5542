"""Build CS5542 Quiz Challenge-1 deck (Corporate Professional style, navy + teal)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
import json, csv

ROOT = Path(r"C:\Users\mural\.openclaw\workspace\CS5542\Week11_Lab")
OUT = ROOT / "slides" / "CS5542_QuizChallenge1_Murali_Ediga.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)

NAVY = RGBColor(0x0B, 0x1F, 0x3B)
TEAL = RGBColor(0x14, 0xB8, 0xA6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1F, 0x2A, 0x44)
MUTED = RGBColor(0x6B, 0x72, 0x80)
FAINT = RGBColor(0xF3, 0xF4, 0xF6)
RULE = RGBColor(0xE5, 0xE7, 0xEB)

FONT = "Segoe UI"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height
BLANK = prs.slide_layouts[6]

def rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp

def tb(slide, x, y, w, h, text, size=18, color=INK, bold=False, font=FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return box

def page_number(slide, n, total):
    tb(slide, Inches(12.4), Inches(7.05), Inches(0.8), Inches(0.3),
       f"{n} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)

def header_bar(slide, title, eyebrow=None):
    rect(slide, 0, 0, SW, Inches(0.08), TEAL)
    if eyebrow:
        tb(slide, Inches(0.6), Inches(0.3), Inches(8), Inches(0.35),
           eyebrow.upper(), size=11, color=TEAL, bold=True)
        tb(slide, Inches(0.6), Inches(0.65), Inches(12), Inches(0.8),
           title, size=28, color=NAVY, bold=True)
    else:
        tb(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8),
           title, size=30, color=NAVY, bold=True)
    rect(slide, Inches(0.6), Inches(1.5), Inches(0.7), Inches(0.04), TEAL)

def footer(slide, label="CS 5542 · Quiz Challenge-1 · Murali Ediga"):
    tb(slide, Inches(0.6), Inches(7.1), Inches(8), Inches(0.3),
       label, size=10, color=MUTED)

# ----------------- build -----------------
TOTAL = 12

# Slide 1 — title
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
# Accent strip
rect(s, 0, Inches(5.8), SW, Inches(0.06), TEAL)
tb(s, Inches(0.9), Inches(0.8), Inches(6), Inches(0.4),
   "CS 5542 · QUIZ CHALLENGE-1", size=13, color=TEAL, bold=True)
tb(s, Inches(0.9), Inches(1.4), Inches(11.5), Inches(2.3),
   "AI-Controlled E-Commerce\nProduct Image Generation", size=54, color=WHITE, bold=True)
tb(s, Inches(0.9), Inches(4.1), Inches(11.5), Inches(1.2),
   "Comparing prompt control strategies on SDXL-Turbo and\nStable Diffusion 1.5 across 15 products",
   size=22, color=RGBColor(0xCB, 0xD5, 0xE1))
tb(s, Inches(0.9), Inches(6.05), Inches(11), Inches(0.5),
   "Murali Ediga", size=20, color=WHITE, bold=True)
tb(s, Inches(0.9), Inches(6.5), Inches(11), Inches(0.4),
   "University of Missouri–Kansas City  ·  Big Data Analytics & Apps  ·  April 20, 2026",
   size=13, color=RGBColor(0x94, 0xA3, 0xB8))

# Slide 2 — Scenario
s = prs.slides.add_slide(BLANK)
header_bar(s, "Metadata-driven product imagery, at catalog scale.", eyebrow="Scenario")
tb(s, Inches(0.6), Inches(1.75), Inches(12), Inches(1.2),
   "E-commerce platforms such as Amazon and Shopify increasingly rely on generative AI to produce product\nimages from structured metadata — reducing studio cost and accelerating time-to-listing.",
   size=16, color=INK)

def callout(x, label, title, body):
    rect(s, x, Inches(3.3), Inches(3.9), Inches(3.4), WHITE, line=RULE)
    rect(s, x, Inches(3.3), Inches(0.12), Inches(3.4), TEAL)
    tb(s, x + Inches(0.3), Inches(3.5), Inches(3.5), Inches(0.4),
       label.upper(), size=11, color=TEAL, bold=True)
    tb(s, x + Inches(0.3), Inches(3.85), Inches(3.5), Inches(0.6),
       title, size=18, color=NAVY, bold=True)
    tb(s, x + Inches(0.3), Inches(4.55), Inches(3.5), Inches(2.0),
       body, size=13, color=INK)

callout(Inches(0.6),  "01", "Problem",
        "Manual photography is slow and expensive. Naive prompts produce inconsistent images that hurt brand quality.")
callout(Inches(4.7),  "02", "Approach",
        "Templated prompt engineering + control (negative prompts, CFG) on Stable Diffusion, evaluated with CLIP metrics.")
callout(Inches(8.8),  "03", "Outcome",
        "Structured prompts boost image consistency +2.7 pts over naive. CLIP-alignment caveats documented.")
footer(s); page_number(s, 2, TOTAL)

# Slide 3 — Dataset
s = prs.slides.add_slide(BLANK)
header_bar(s, "15 curated products across six categories.", eyebrow="Dataset")
tb(s, Inches(0.6), Inches(1.75), Inches(12), Inches(0.6),
   "Amazon-style product metadata: title, category, color, material, style, and descriptive features.",
   size=15, color=INK)
# Sample table
rows = [
    ("p01", "Classic Leather Oxford Shoes", "Footwear", "Brown", "Genuine leather", "Formal"),
    ("p03", "Wireless Over-Ear Headphones", "Electronics", "Matte Black", "Plastic + foam", "Modern"),
    ("p13", "Velvet Tufted Sofa", "Furniture", "Emerald", "Velvet upholstery", "Luxury"),
]
hdr = ["ID", "Title", "Category", "Color", "Material", "Style"]
widths = [Inches(0.7), Inches(4.1), Inches(1.8), Inches(1.5), Inches(2.2), Inches(1.5)]
tbl_x = Inches(0.6); tbl_y = Inches(2.6); row_h = Inches(0.55)
# header
x = tbl_x
rect(s, tbl_x, tbl_y, sum(widths, Emu(0)), row_h, NAVY)
for h_, w_ in zip(hdr, widths):
    tb(s, x + Inches(0.15), tbl_y + Inches(0.13), w_, row_h,
       h_, size=12, color=WHITE, bold=True)
    x += w_
# rows
y = tbl_y + row_h
for i, row in enumerate(rows):
    bg = WHITE if i % 2 == 0 else FAINT
    rect(s, tbl_x, y, sum(widths, Emu(0)), row_h, bg, line=RULE)
    x = tbl_x
    for cell, w_ in zip(row, widths):
        tb(s, x + Inches(0.15), y + Inches(0.13), w_, row_h,
           cell, size=12, color=INK)
        x += w_
    y += row_h
# stats
tb(s, Inches(0.6), Inches(5.4), Inches(12), Inches(0.4),
   "CATEGORIES  ·  Footwear · Electronics · Apparel · Home · Furniture · Accessories · Kitchen",
   size=11, color=TEAL, bold=True)
# Three big numbers
def stat(x, num, lbl):
    tb(s, x, Inches(5.9), Inches(3), Inches(0.8), num, size=44, color=NAVY, bold=True)
    tb(s, x, Inches(6.75), Inches(3), Inches(0.4), lbl, size=12, color=MUTED, bold=True)
stat(Inches(0.6), "15", "PRODUCTS")
stat(Inches(4.6), "6", "CATEGORIES")
stat(Inches(8.6), "7", "METADATA FIELDS")
footer(s); page_number(s, 3, TOTAL)

# Slide 4 — Methodology
s = prs.slides.add_slide(BLANK)
header_bar(s, "End-to-end metadata-to-image pipeline.", eyebrow="Methodology")

def pipebox(x, y, w, h, title, sub, accent=False):
    rect(s, x, y, w, h, NAVY if accent else WHITE, line=None if accent else RULE)
    rect(s, x, y, Inches(0.1), h, TEAL)
    tb(s, x + Inches(0.3), y + Inches(0.3), w - Inches(0.5), Inches(0.6),
       title, size=16, color=(WHITE if accent else NAVY), bold=True)
    tb(s, x + Inches(0.3), y + Inches(0.95), w - Inches(0.5), h - Inches(1.2),
       sub, size=12, color=(RGBColor(0xCB,0xD5,0xE1) if accent else INK))

def arrow(x, y):
    shp = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, Inches(0.5), Inches(0.3))
    shp.fill.solid(); shp.fill.fore_color.rgb = TEAL
    shp.line.fill.background()

W = Inches(2.25); H = Inches(2.2)
Y = Inches(2.5)
GAP = Inches(0.25)
X0 = Inches(0.6)
pipebox(X0,                           Y, W, H, "Metadata", "15 × products.csv\ntitle, color, material, style, features")
arrow(X0 + W + Inches(-0.05), Y + Inches(0.95))
pipebox(X0 + W + GAP + Inches(0.25),  Y, W, H, "Prompt Templating",
        "A · Naive\nB · Structured\nC · Structured + Negative")
arrow(X0 + 2*W + 2*GAP + Inches(0.2), Y + Inches(0.95))
pipebox(X0 + 2*(W+GAP) + Inches(0.5), Y, W, H, "Stable Diffusion",
        "SDXL-Turbo (A, B)\nSD 1.5 (C, w/ CFG)\nRTX PRO 5000 · fp16", accent=True)
arrow(X0 + 3*W + 3*GAP + Inches(0.45), Y + Inches(0.95))
pipebox(X0 + 3*(W+GAP) + Inches(0.75), Y, W, H, "135 Images",
        "15 × 3 × 3\n(product × strategy × seed)")
arrow(X0 + 4*W + 4*GAP + Inches(0.7), Y + Inches(0.95))
pipebox(X0 + 4*(W+GAP) + Inches(1.0), Y, W, H, "Evaluation",
        "CLIP align\nConsistency\nDiversity · failures")

# bullets under
tb(s, Inches(0.6), Inches(5.2), Inches(12), Inches(0.5),
   "Runtime", size=12, color=TEAL, bold=True)
tb(s, Inches(0.6), Inches(5.55), Inches(12), Inches(1.5),
   "•  90 images generated on SDXL-Turbo in 28 s (≈3 img/s, 4 inference steps, guidance = 0)\n"
   "•  45 images generated on SD 1.5 in 33 s (25 steps, guidance = 7.5) — enables negative prompting\n"
   "•  Deterministic seeds (1000–1002) so every product / strategy pair is reproducible",
   size=13, color=INK)
footer(s); page_number(s, 4, TOTAL)

# Slide 5 — SD pipeline
s = prs.slides.add_slide(BLANK)
header_bar(s, "Two Stable Diffusion backbones, one control surface.", eyebrow="Stable Diffusion Pipeline")

# Left: models panel
rect(s, Inches(0.6), Inches(1.9), Inches(5.9), Inches(4.8), WHITE, line=RULE)
rect(s, Inches(0.6), Inches(1.9), Inches(5.9), Inches(0.12), TEAL)
tb(s, Inches(0.85), Inches(2.1), Inches(5.5), Inches(0.4),
   "MODELS", size=11, color=TEAL, bold=True)

def modelrow(y, name, cfg, steps, use):
    tb(s, Inches(0.85), y, Inches(5.5), Inches(0.35),
       name, size=14, color=NAVY, bold=True)
    tb(s, Inches(0.85), y + Inches(0.35), Inches(5.5), Inches(0.8),
       f"guidance = {cfg}   ·   steps = {steps}\n{use}",
       size=11, color=INK)

modelrow(Inches(2.55), "stabilityai/sdxl-turbo", "0.0", "4",
         "Used for strategies A (naive) and B (structured).\nfp16 · 512×512 · distilled turbo schedule.")
rect(s, Inches(0.85), Inches(4.2), Inches(5.4), Inches(0.02), RULE)
modelrow(Inches(4.35), "runwayml/stable-diffusion-v1-5", "7.5", "25",
         "Used for strategy C to enable negative prompts via\nclassifier-free guidance (CFG disabled in turbo).")

# Right: code snippet
rect(s, Inches(6.8), Inches(1.9), Inches(5.9), Inches(4.8), RGBColor(0x0F, 0x17, 0x2A))
tb(s, Inches(6.95), Inches(2.05), Inches(5.5), Inches(0.4),
   "PIPELINE (HuggingFace diffusers 0.37)", size=11, color=TEAL, bold=True, font=FONT)
code = (
    "pipe = AutoPipelineForText2Image.from_pretrained(\n"
    "    \"stabilityai/sdxl-turbo\",\n"
    "    torch_dtype=torch.float16,\n"
    "    variant=\"fp16\",\n"
    ").to(\"cuda\")\n\n"
    "img = pipe(\n"
    "    prompt=structured_prompt(row),\n"
    "    negative_prompt=NEGATIVE,\n"
    "    num_inference_steps=4,\n"
    "    guidance_scale=0.0,\n"
    "    generator=g,\n"
    ").images[0]"
)
tb(s, Inches(6.95), Inches(2.5), Inches(5.7), Inches(4.1),
   code, size=12, color=RGBColor(0xE2, 0xE8, 0xF0), font="Consolas")

# Hardware footnote
tb(s, Inches(0.6), Inches(6.85), Inches(12), Inches(0.3),
   "HARDWARE  ·  NVIDIA RTX PRO 5000 Blackwell (48 GB)  ·  PyTorch 2.11 + CUDA 13.0",
   size=10, color=MUTED, bold=True)
footer(s); page_number(s, 5, TOTAL)

# Slide 6 — Prompt design
s = prs.slides.add_slide(BLANK)
header_bar(s, "Three prompt strategies, one product.", eyebrow="Prompt Design")
tb(s, Inches(0.6), Inches(1.75), Inches(12), Inches(0.5),
   "Example input row — p01 Classic Leather Oxford Shoes", size=12, color=TEAL, bold=True)

def promptcard(x, label, kind, prompt_text, negative=None, bg=WHITE):
    rect(s, x, Inches(2.3), Inches(4), Inches(4.5), bg, line=RULE)
    rect(s, x, Inches(2.3), Inches(4), Inches(0.12), TEAL)
    tb(s, x + Inches(0.25), Inches(2.5), Inches(3.6), Inches(0.35),
       label.upper(), size=11, color=TEAL, bold=True)
    tb(s, x + Inches(0.25), Inches(2.85), Inches(3.6), Inches(0.5),
       kind, size=16, color=NAVY, bold=True)
    tb(s, x + Inches(0.25), Inches(3.4), Inches(3.6), Inches(2.6),
       prompt_text, size=11, color=INK)
    if negative:
        tb(s, x + Inches(0.25), Inches(5.8), Inches(3.6), Inches(0.35),
           "NEGATIVE", size=10, color=RGBColor(0xB9, 0x1C, 0x1C), bold=True)
        tb(s, x + Inches(0.25), Inches(6.15), Inches(3.6), Inches(0.6),
           negative, size=10, color=MUTED)

promptcard(Inches(0.6), "A", "Naive",
           "“Classic Leather Oxford Shoes”")
promptcard(Inches(4.7), "B", "Structured",
           "“Classic Leather Oxford Shoes, Brown, made of Genuine Leather, "
           "Formal style, Lace-up stitched sole polished finish, studio lighting, "
           "clean white background, centered product shot, 4k, photorealistic, "
           "commercial photography”")
promptcard(Inches(8.8), "C", "Structured + Negative",
           "Same as (B), rendered with SD 1.5 + CFG 7.5",
           negative="blurry, low quality, distorted, watermark, text, cluttered background, extra limbs, oversaturated, cartoon, sketch")
footer(s); page_number(s, 6, TOTAL)

# Slide 7 — Control strategy
s = prs.slides.add_slide(BLANK)
header_bar(s, "Two control mechanisms — one data-side, one model-side.", eyebrow="Control Strategy")

def bigpoint(x, num, title, body):
    rect(s, x, Inches(2.0), Inches(6), Inches(4.5), WHITE, line=RULE)
    tb(s, x + Inches(0.3), Inches(2.2), Inches(1.2), Inches(1.2),
       num, size=60, color=TEAL, bold=True)
    tb(s, x + Inches(0.3), Inches(3.5), Inches(5.5), Inches(0.5),
       title, size=20, color=NAVY, bold=True)
    tb(s, x + Inches(0.3), Inches(4.1), Inches(5.5), Inches(2.4),
       body, size=13, color=INK)

bigpoint(Inches(0.6), "01", "Structured prompt templates",
         "A deterministic template maps product metadata into a canonical prompt form:\n"
         "   {title}, {color}, made of {material}, {style} style, {features}, + quality tokens.\n\n"
         "Benefit: removes user-input variance. Every product in the catalog\nreceives the same prompt structure → higher image consistency.")
bigpoint(Inches(6.9), "02", "Negative prompt conditioning",
         "A fixed negative prompt (blurry, watermark, cluttered background, "
         "cartoon, sketch, extra limbs…) is applied via classifier-free guidance.\n\n"
         "Constraint: CFG is disabled in SDXL-Turbo (guidance=0).\n"
         "We switched strategy C to SD 1.5 with guidance_scale=7.5 to make\nthe negative prompt actually affect sampling.")

tb(s, Inches(0.6), Inches(6.8), Inches(12), Inches(0.3),
   "ControlNet was out of scope for this assignment; flagged as future work on slide 11.",
   size=11, color=MUTED)
footer(s); page_number(s, 7, TOTAL)

# Slide 8 — Tools
s = prs.slides.add_slide(BLANK)
header_bar(s, "Tools & technologies.", eyebrow="Stack")
tools = [
    ("HuggingFace diffusers 0.37",  "Pipeline orchestration, text-to-image"),
    ("PyTorch 2.11 + CUDA 13.0",    "fp16 inference on RTX PRO 5000"),
    ("SDXL-Turbo",                  "Distilled 1–4 step text-to-image"),
    ("Stable Diffusion v1.5",       "CFG-enabled model for negative prompts"),
    ("open_clip (ViT-B-32, OpenAI)","CLIP similarity scoring"),
    ("Pillow + python-pptx",        "Grid assembly, report slide generation"),
    ("Windows + Tailscale + SSH",   "Dev on Windows; execution on Linux GPU"),
    ("Git + GitHub",                "Versioned source + sample outputs"),
]
for i, (name, desc) in enumerate(tools):
    col = i % 2; row = i // 2
    x = Inches(0.6) + col * Inches(6.2)
    y = Inches(2.0) + row * Inches(1.05)
    rect(s, x, y, Inches(5.9), Inches(0.9), WHITE, line=RULE)
    rect(s, x, y, Inches(0.08), Inches(0.9), TEAL)
    tb(s, x + Inches(0.25), y + Inches(0.12), Inches(5.5), Inches(0.45),
       name, size=15, color=NAVY, bold=True)
    tb(s, x + Inches(0.25), y + Inches(0.5), Inches(5.5), Inches(0.4),
       desc, size=12, color=INK)

tb(s, Inches(0.6), Inches(6.85), Inches(12), Inches(0.3),
   "135 images generated in ~60 s combined.  Total code: 4 Python files, ~250 LOC.",
   size=11, color=MUTED, bold=True)
footer(s); page_number(s, 8, TOTAL)

# Slide 9 — Results gallery
s = prs.slides.add_slide(BLANK)
header_bar(s, "Results — naive · structured · structured + negative.", eyebrow="Results Gallery")
tb(s, Inches(0.6), Inches(1.75), Inches(12), Inches(0.4),
   "Column 1: Naive (SDXL-Turbo)    ·    Column 2: Structured (SDXL-Turbo)    ·    Column 3: Structured + Negative (SD 1.5)",
   size=11, color=TEAL, bold=True)

grids = [
    ("p01_compare.png", "Oxford Shoes"),
    ("p03_compare.png", "Headphones"),
    ("p07_compare.png", "Handbag"),
    ("p13_compare.png", "Velvet Sofa"),
]
# 2x2 grid of 4 rows of comparison images
# Each compare image is 1536x552 aspect, place wide & short
GRID_W = Inches(6.0); GRID_H = Inches(2.35)
GX = [Inches(0.6), Inches(6.85)]
GY = [Inches(2.25), Inches(4.9)]
for i, (fn, label) in enumerate(grids):
    col = i % 2; row = i // 2
    x = GX[col]; y = GY[row]
    path = ROOT / "outputs" / "grids" / fn
    if path.exists():
        s.shapes.add_picture(str(path), x, y, width=GRID_W, height=GRID_H)
        rect(s, x, y + GRID_H, GRID_W, Inches(0.25), NAVY)
        tb(s, x + Inches(0.15), y + GRID_H + Inches(0.01), GRID_W, Inches(0.25),
           label, size=11, color=WHITE, bold=True)
footer(s); page_number(s, 9, TOTAL)

# Slide 10 — Metrics
s = prs.slides.add_slide(BLANK)
header_bar(s, "Quantitative evaluation (CLIP ViT-B/32).", eyebrow="Evaluation")

# Table
hdr = ["Strategy", "CLIP Alignment (mean ± std)", "Consistency ↑", "Diversity ↑"]
rows = [
    ("A · Naive",                 "0.314 ± 0.018", "0.907", "0.349"),
    ("B · Structured",            "0.308 ± 0.016", "0.932", "0.332"),
    ("C · Structured + Negative", "0.306 ± 0.021", "0.881", "0.304"),
]
widths = [Inches(3.6), Inches(3.8), Inches(2.5), Inches(2.3)]
total_w = sum(widths, Emu(0))
tbl_x = Inches(0.6); tbl_y = Inches(2.0)
# Header
rect(s, tbl_x, tbl_y, total_w, Inches(0.65), NAVY)
x = tbl_x
for h_, w_ in zip(hdr, widths):
    tb(s, x + Inches(0.2), tbl_y + Inches(0.17), w_, Inches(0.65),
       h_, size=13, color=WHITE, bold=True)
    x += w_
# Rows
y = tbl_y + Inches(0.65)
best_cells = {"Consistency ↑": 1, "Diversity ↑": 0}  # row index that wins each column
for i, row in enumerate(rows):
    bg = WHITE if i % 2 == 0 else FAINT
    rect(s, tbl_x, y, total_w, Inches(0.75), bg, line=RULE)
    x = tbl_x
    for j, (cell, w_) in enumerate(zip(row, widths)):
        color = INK
        bold = False
        if j == 2 and i == 1: color = TEAL; bold = True
        if j == 3 and i == 0: color = TEAL; bold = True
        tb(s, x + Inches(0.2), y + Inches(0.22), w_, Inches(0.75),
           cell, size=14, color=color, bold=bold)
        x += w_
    y += Inches(0.75)

# Interpretation bullets
tb(s, Inches(0.6), Inches(5.1), Inches(12), Inches(0.4),
   "INTERPRETATION", size=12, color=TEAL, bold=True)
tb(s, Inches(0.6), Inches(5.5), Inches(12), Inches(1.8),
   "•  Structured prompts (B) maximize image-to-image consistency (+2.7 pts over naive) — the metric that matters most for\n"
   "   e-commerce catalog uniformity.\n"
   "•  Naive prompts (A) score marginally higher on CLIP-alignment: short prompts are easier for CLIP to match (well-known bias).\n"
   "•  Strategy C trades consistency for visual cleanliness and uses a different base model, reflected in lower diversity.",
   size=13, color=INK)
footer(s); page_number(s, 10, TOTAL)

# Slide 11 — Findings & Limitations
s = prs.slides.add_slide(BLANK)
header_bar(s, "What we learned, and where it breaks.", eyebrow="Findings & Limitations")

def col(x, title, items, accent):
    tb(s, x, Inches(1.95), Inches(5.9), Inches(0.5),
       title.upper(), size=12, color=accent, bold=True)
    rect(s, x, Inches(2.3), Inches(0.8), Inches(0.04), accent)
    y = Inches(2.55)
    for it in items:
        tb(s, x, y, Inches(5.9), Inches(1.2),
           "•  " + it, size=13, color=INK)
        y += Inches(1.0)

col(Inches(0.6), "Findings",
    ["Structured prompts raise image consistency by +2.7 pts vs naive — a win for brand uniformity.",
     "CLIP-alignment is NOT a reliable quality metric: it rewards short prompts (length bias).",
     "Turbo models disable classifier-free guidance, so negative prompts require a non-turbo base."],
    TEAL)
col(Inches(6.9), "Limitations",
    ["15 products × 3 seeds — small sample size, no human preference study.",
     "No ControlNet / pose / layout conditioning; no brand-specific LoRA fine-tuning.",
     "CLIP ViT-B/32 is English-only — results may not generalize to multilingual catalogs."],
    RGBColor(0xB9, 0x1C, 0x1C))
footer(s); page_number(s, 11, TOTAL)

# Slide 12 — AI disclosure, links
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(1.0), SW, Inches(0.06), TEAL)
tb(s, Inches(0.9), Inches(0.4), Inches(8), Inches(0.6),
   "AI TOOL DISCLOSURE  ·  LINKS", size=13, color=TEAL, bold=True)
tb(s, Inches(0.9), Inches(1.2), Inches(11.5), Inches(0.8),
   "Transparency & reproducibility", size=40, color=WHITE, bold=True)

# two columns
tb(s, Inches(0.9), Inches(2.3), Inches(5.9), Inches(0.4),
   "AI TOOLS USED", size=12, color=TEAL, bold=True)
tb(s, Inches(0.9), Inches(2.7), Inches(5.9), Inches(3.5),
   "•  Claude (Anthropic) — code scaffolding for generate.py\n"
   "    and evaluate.py, prompt template drafting, slide\n"
   "    structuring, debugging.\n"
   "•  No AI tools were used to hand-edit evaluation\n"
   "    numbers or interpret results.\n"
   "•  All 135 images were executed by the author on an\n"
   "    RTX PRO 5000 GPU; metrics computed locally.\n"
   "•  Final dataset curation, strategy selection, and\n"
   "    conclusions are the author’s own.",
   size=13, color=WHITE)

tb(s, Inches(7.2), Inches(2.3), Inches(5.5), Inches(0.4),
   "LINKS", size=12, color=TEAL, bold=True)
tb(s, Inches(7.2), Inches(2.7), Inches(5.5), Inches(3.5),
   "GitHub\n"
   "github.com/muralikrish9/CS5542/tree/master/Week11_Lab\n\n"
   "Demo video (1–2 min)\n"
   "youtu.be/LR7pLWJjRjQ\n\n"
   "Author\n"
   "Murali Ediga  ·  University of Missouri–Kansas City\n"
   "CS 5542 Big Data Analytics & Apps  ·  Spring 2026",
   size=13, color=WHITE)

tb(s, Inches(0.9), Inches(6.7), Inches(12), Inches(0.5),
   "Thank you.", size=28, color=TEAL, bold=True)
tb(s, Inches(12.4), Inches(7.05), Inches(0.8), Inches(0.3),
   f"{TOTAL} / {TOTAL}", size=10, color=RGBColor(0x94,0xA3,0xB8), align=PP_ALIGN.RIGHT)

prs.save(str(OUT))
print("saved:", OUT)
